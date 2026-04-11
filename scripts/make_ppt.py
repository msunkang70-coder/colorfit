"""ColorFit 멘토 발표 PPT 생성."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

MARSALA = RGBColor(0x96, 0x4F, 0x4C)
BG = RGBColor(0xF8, 0xF6, 0xF3)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x8C, 0x85, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PINK = RGBColor(0xFF, 0xDD, 0xDA)


def make_slide(title, subtitle, bullets, note, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG

    # number
    tb = s.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{num}/12"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

    # accent bar
    sh = s.shapes.add_shape(1, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3))
    sh.fill.solid()
    sh.fill.fore_color.rgb = MARSALA
    sh.line.fill.background()

    # title
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    # subtitle
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(16)
    p.font.color.rgb = MARSALA
    p.font.italic = True

    # bullets
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)

    # speaker notes
    if note:
        s.notes_slide.notes_text_frame.text = note


def make_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = MARSALA

    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11), Inches(3.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ColorFit"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "결정을 쉽게 만드는 스타일 서비스"
    p2.font.size = Pt(20)
    p2.font.color.rgb = PINK

    p3 = tf.add_paragraph()
    p3.text = ""
    p4 = tf.add_paragraph()
    p4.text = "멘토 발표 | 2026.04"
    p4.font.size = Pt(14)
    p4.font.color.rgb = PINK


def make_qa():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = MARSALA

    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(2.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Q & A"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = ""
    p3 = tf.add_paragraph()
    p3.text = "ColorFit \u2014 \uacb0\uc815\uc744 \uc27d\uac8c \ub9cc\ub4dc\ub294 \uc2a4\ud0c0\uc77c \uc11c\ube44\uc2a4"
    p3.font.size = Pt(18)
    p3.font.color.rgb = PINK

    qa_text = (
        "Q1: \uc65c \ube0c\ub79c\uce58 \ub098\ub220\ub098\uc694? \u2192 \uac19\uc740 \uc11c\ube44\uc2a4\uc5d0\uc11c \uac01\uc790 \ub2e4\ub978 \uad00\uc810\uc73c\ub85c \uc2e4\ud5d8. \ubcd1\ub82c \uac80\uc99d.\n"
        "Q2: \uc65c \ud1b5\ud569 \uc548\ud588\ub098\uc694? \u2192 \uac80\uc99d \ub2e8\uacc4. \ud6a8\uacfc \ud655\uc778 \ud6c4 \ud1b5\ud569 \uc608\uc815.\n"
        "Q3: \uae30\uc874 \uc11c\ube44\uc2a4\uc640 \ucc28\uc774? \u2192 10\uac1c \ucd94\ucc9c vs 1\uac1c \uacb0\uc815. \ud575\uc2ec\uc740 \uacb0\uc815 \uad6c\uc870.\n"
        "Q4: \uac1c\uc778\ud654 \uc218\uc900? \u2192 \ud1a4+TPO+\uc608\uc0b0+\ucde8\ud5a5 4\uac00\uc9c0. \ud589\ub3d9 \ub370\uc774\ud130 \uc218\uc9d1 \uc911.\n"
        "Q5: DB \uc65c \uc548 \uc37c\ub098\uc694? \u2192 MVP \uc18d\ub3c4 \uc6b0\uc120. \ubaa8\ub378 \uc124\uacc4 \uc644\ub8cc, PostgreSQL \uc804\ud658 \uc900\ube44\ub428.\n"
        "Q6: \ud655\uc7a5 \uac00\ub2a5\uc131? \u2192 DB \uc804\ud658/\ud589\ub3d9 \ud559\uc2b5/\uc804\ubb38\uac00 \uaddc\uce59 3\ubc29\ud5a5 \ubaa8\ub450 \uac00\ub2a5.\n"
        "Q7: \uc218\uc775 \ubaa8\ub378? \u2192 CTA \uae30\ubc18 \uc81c\ud734 \ucee4\uba38\uc2a4(CPA). Decision \uc11c\ube44\uc2a4\uc758 \uc804\ud658\uc728 \uac15\uc810.\n"
        "Q8: \uac00\uc7a5 \ud070 \ub9ac\uc2a4\ud06c? \u2192 \uc804\ubb38\uac00 \uc124\ub4dd\ub825 \ubd80\uc871 + \ub370\uc774\ud130 \ub2e4\uc591\uc131 \ud55c\uacc4."
    )
    s.notes_slide.notes_text_frame.text = qa_text


# ═══ BUILD ═══

make_cover()

SLIDES = [
    (
        "\ucd94\ucc9c\uc774 \ub9ce\uc744\uc218\ub85d \uacb0\uc815\uc740 \ub354 \uc5b4\ub824\uc6cc\uc9c4\ub2e4",
        '"\ucd94\ucc9c 10\uac1c\ub97c \ubcf4\uc5ec\uc8fc\uba74 \uc0ac\uc6a9\uc790\ub294 \uace0\ub974\ub294 \uac8c \uc544\ub2c8\ub77c \ud3ec\uae30\ud55c\ub2e4"',
        [
            "\u2022 \ud328\uc158 \ucd94\ucc9c \uc11c\ube44\uc2a4 \uc774\ud0c8\ub960 60%+ \u2014 \ub300\ubd80\ubd84 \"\uacb0\uacfc \ud654\uba74\"\uc5d0\uc11c \uc774\ud0c8",
            "\u2022 \uc6d0\uc778: \uc120\ud0dd\uc9c0 \uacfc\uc789 (Choice Overload)",
            "\u2022 10\uac1c \ucd94\ucc9c = 10\ubc88\uc758 \ube44\uad50 \ud310\ub2e8 \u2192 \ud53c\ub85c \u2192 \uc774\ud0c8",
            '\u2022 \uc0ac\uc6a9\uc790\uac00 \uc6d0\ud558\ub294 \uac74 "\uace8\ub77c\uc918"',
        ],
        "\ud328\uc158 \ucd94\ucc9c \uc11c\ube44\uc2a4\ub97c \uc4f0\uba74\uc11c \uacb0\uacfc \ud654\uba74\uc5d0\uc11c \uc774\ud0c8\ud55c \uacbd\ud5d8 \ub2e4\ub4e4 \uc788\uc73c\uc2e4 \uac81\ub2c8\ub2e4. \uc800\ud76c \ud300\uc774 \uc8fc\ubaa9\ud55c \uac74 \ucd94\ucc9c\uc740 \uc798 \ub418\ub294\ub370 \uacb0\uc815\uc744 \ubabb \ud55c\ub2e4\ub294 \ubb38\uc81c\uc785\ub2c8\ub2e4.",
    ),
    (
        "\ucd94\ucc9c \uc11c\ube44\uc2a4\uac00 \uc544\ub2c8\ub77c \uacb0\uc815 \uc11c\ube44\uc2a4",
        '"\uc6b0\ub9ac\ub294 Top1\uc744 \uacb0\uc815\ud558\uac8c \ub9cc\ub4dc\ub294 \uad6c\uc870\ub97c \uc124\uacc4\ud588\ub2e4"',
        [
            '\u2022 \ucd94\ucc9c \uc11c\ube44\uc2a4: "\uc5ec\uae30\uc11c \uace8\ub77c" \u2192 \uc120\ud0dd \ubd80\ub2f4 \uc804\uac00',
            '\u2022 Decision \uc11c\ube44\uc2a4: "\uc774\uac78\ub85c \uac00" \u2192 \ud310\ub2e8 \uadfc\uac70 \uc81c\uacf5',
            "\u2022 \ud750\ub984: Top1 \uc81c\uc2dc \u2192 \uc65c \uc774\uac74\uc9c0 \uc124\uba85 \u2192 \ud655\uc2e0 \u2192 \uacb0\uc815",
            "\u2022 \ud544\uc694\ud560 \ub54c\ub9cc \ube44\uad50 (Explore Mode)",
        ],
        "ColorFit\uc740 \ucd94\ucc9c \uc11c\ube44\uc2a4\uac00 \uc544\ub2c8\ub77c \uacb0\uc815 \uc11c\ube44\uc2a4\uc785\ub2c8\ub2e4. \ud558\ub098\ub97c \uc81c\uc2dc\ud558\uace0, \uc65c \uc774\uac74\uc9c0 \uc124\uba85\ud558\uace0, \ud655\uc2e0\uc774 \uc0dd\uae30\uba74 \ubc14\ub85c \ud589\ub3d9\ud558\uac8c \ub9cc\ub4ed\ub2c8\ub2e4.",
    ),
    (
        "\uacf5\ud1b5 \uad6c\uc870 \uc704\uc5d0\uc11c \uac01\uc790 \uc2e4\ud5d8\ud558\ub294 \ud611\uc5c5",
        '"\uac19\uc740 \uc11c\ube44\uc2a4, \ub2e4\ub978 \uad00\uc810\uc73c\ub85c \uac80\uc99d"',
        [
            "\u2022 Main branch: PM\uc774 \uad00\ub9ac\ud558\ub294 \uacf5\ud1b5 \uad6c\uc870 (\ucd94\ucc9c \ud30c\uc774\ud504\ub77c\uc778, \ub370\uc774\ud130, API)",
            '\u2022 \uac1c\uc778 branch: \uac01\uc790\uc758 \uad00\uc810\uc73c\ub85c \uac1c\uc120 \ubc29\ud5a5 \uc2e4\ud5d8',
            '\u2022 \uc800\ub294 "Decision UX \uc124\ub4dd\ub825 \uac15\ud654" \ubc29\ud5a5\uc73c\ub85c \uc791\uc5c5',
            "\u2022 \ud1b5\ud569\uc774 \uc544\ub2c8\ub77c \ubcd1\ub82c \uac80\uc99d \u2192 \ub354 \ub098\uc740 \ubc29\ud5a5 \ub3c4\ucd9c",
        ],
        "\uc800\ud76c \ud300\uc740 \ud558\ub098\uc758 \uc11c\ube44\uc2a4\ub97c \uacf5\ub3d9\uc73c\ub85c \ub9cc\ub4e4\ub418, \uac01\uc790 \ub2e4\ub978 \ubc29\ud5a5\uc5d0\uc11c \uac1c\uc120\uc744 \uc2e4\ud5d8\ud558\ub294 \uad6c\uc870\uc785\ub2c8\ub2e4.",
    ),
    (
        "Decision UX\ub97c \uc704\ud55c \ucd94\ucc9c \ud30c\uc774\ud504\ub77c\uc778",
        '"\uc774\ubbf8 \uad6c\uc870\ub294 Decision-First\ub85c \uc124\uacc4\ub418\uc5b4 \uc788\ub2e4"',
        [
            "\u2022 \uc785\ub825: \ud37c\uc2a4\ub110\ucee8\ub7ec(12\ud1a4) + TPO(8\uc885) + \uc608\uc0b0 + \uc2a4\ud0c0\uc77c \ucde8\ud5a5",
            "\u2022 Hard Filter: \uc131\ubcc4\xb7\uc608\uc0b0\xb7\uc2dc\uc98c \ud0c8\ub77d \u2192 Soft Score: 5\ucd95 \uc810\uc218 \uc815\ub82c",
            "\u2022 5\ucd95: TPO \uc801\ud569\ub3c4 / \ud54f / \ucee8\ub7ec / \uc2a4\ud0c0\uc77c \uc77c\uad00\uc131 / \ub9ac\uc2a4\ud06c",
            "\u2022 \ucd9c\ub825: Top1 (Decision Mode) + Top3 (Explore Mode)",
        ],
        "\uc0ac\uc6a9\uc790\uac00 \uc870\uac74\uc744 \uc785\ub825\ud558\uba74 Hard Filter\ub85c \ubd80\uc801\ud569\ud55c \ucf54\ub514\ub97c \uc81c\uac70\ud558\uace0, 5\ucd95 \uc810\uc218\ub85c Top1\uc744 \uacb0\uc815\ud569\ub2c8\ub2e4.",
    ),
    (
        "\ucd94\ucc9c \uc815\ud655\ub3c4\uac00 \uc544\ub2c8\ub77c \uacb0\uc815 \ud655\ub960\uc5d0 \uc9d1\uc911",
        '"\uc0ac\uc6a9\uc790\uac00 Top1\uc744 \ubcf4\uace0 \uc774\uac78\ub85c \uac00\uc57c\uaca0\ub2e4\uace0 \ub290\ub07c\uac8c \ub9cc\ub4dc\ub294 \uc791\uc5c5"',
        [
            "\u2022 Top1 \uc124\ub4dd\ub825 \uac15\ud654: TPO\ubcc4 \uc804\ubb38\uac00 \ud310\ub2e8 \uae30\uc900 (Stylist Criteria Layer)",
            '\u2022 \ub9ac\uc2a4\ud06c \uae30\ubc18 \uc124\uba85: "\uc65c \uc548\uc804\ud55c \uc120\ud0dd\uc778\uc9c0" 3\ub2e8 \uad6c\uc870\ub85c \uc81c\uc2dc',
            '\u2022 Explore UX: "\ube44\uc2b7\ud55c \uc120\ud0dd \ubcf4\uae30" \u2192 "\ub2e4\ub978 \uc778\uc0c1\uc758 \uba74\uc811\ub8e9 \ube44\uad50\ud558\uae30"',
            "\u2022 QA \uac80\uc99d \ub3c4\uad6c: \ucd94\ucc9c \uacb0\uacfc\ub97c \uc870\uac74\ubcc4\ub85c \ube44\uad50 \uac80\uc99d\ud558\ub294 \uc2dc\ubbac\ub808\uc774\ud130",
        ],
        "\uc81c \uc791\uc5c5\uc740 \ucd94\ucc9c \uc54c\uace0\ub9ac\uc998\uc744 \ubc14\uafb8\ub294 \uac8c \uc544\ub2c8\ub77c, \uac19\uc740 \ucd94\ucc9c \uacb0\uacfc\ub97c \ub354 \uc804\ubb38\uac00\ub2f5\uac8c \ud574\uc11d\ud558\uace0 \ud45c\ud604\ud558\ub294 \ub370 \uc9d1\uc911\ud588\uc2b5\ub2c8\ub2e4.",
    ),
    (
        "\uc2e4\uc81c \ub3d9\uc791\ud558\ub294 Decision \ud750\ub984",
        '"\uc628\ubcf4\ub529\ubd80\ud130 \uacb0\uc815\uae4c\uc9c0 \uc804\uccb4 \ud750\ub984\uc774 \ub3d9\uc791\ud55c\ub2e4"',
        [
            "\u2022 \uc628\ubcf4\ub529 5\ub2e8\uacc4: \uc131\ubcc4 \u2192 \ud37c\uc2a4\ub110\ucee8\ub7ec \u2192 TPO+\ubb34\ub4dc \u2192 \uc608\uc0b0 \u2192 \uc2a4\ud0c0\uc77c",
            "\u2022 Decision Mode: Top1 \uce74\ub4dc + \uc804\ubb38\uac00 \ucf54\uba58\ud2b8 + \ub9ac\uc2a4\ud06c \uac00\ub4dc",
            "\u2022 CTA \u2192 \uc124\ubb38(\uc2e0\ub8b0\ub3c4\xb7\ud655\uc2e0) \u2192 \uc0c1\ud488 \uc774\ub3d9",
            "\u2022 \uce21\uc815: TTD(\uacb0\uc815 \uc18c\uc694 \uc2dc\uac04), \uc2e0\ub8b0\ub3c4, \ud655\uc2e0\ub3c4",
        ],
        "\ud604\uc7ac \uc628\ubcf4\ub529\ubd80\ud130 \uacb0\uc815, \uce21\uc815\uae4c\uc9c0 \uc804\uccb4 \ud750\ub984\uc774 \uc2e4\uc81c\ub85c \ub3d9\uc791\ud569\ub2c8\ub2e4.",
    ),
    (
        "\ub370\ubaa8: \uc0ac\uc6a9\uc790\uac00 \uace0\ubbfc\ud558\uc9c0 \uc54a\uac8c \ub9cc\ub4dc\ub294 \ud750\ub984",
        '"\uc9c1\uc811 \ubcf4\uc2dc\uba74 Decision \uc11c\ube44\uc2a4\uc758 \ucc28\uc774\uac00 \ub290\uaef4\uc9d1\ub2c8\ub2e4"',
        [
            "\u2460 \uc628\ubcf4\ub529: \uc5ec\uc131 \u2192 \ubd04\uc6dc\ub77c\uc774\ud2b8 \u2192 \uba74\uc811 \u2192 \ubbf8\ub2c8\uba40",
            "\u2461 Top1: \uc804\ubb38\uac00 \ucf54\uba58\ud2b8 + \ud310\ub2e8 \uae30\uc900 \ubc30\uc9c0 (\ub2e8\uc815\ud568 \xb7 \uc548\uc815\uac10 \xb7 \uc2e0\ub8b0)",
            "\u2462 \ub9ac\uc2a4\ud06c \uac00\ub4dc: [Shield] \uaca9\uc2dd \uc77c\uad00 \u2192 \uc2e4\ud328 \uac00\ub2a5\uc131 \ub0ae\uc74c",
            '\u2463 Explore: "\ub2e4\ub978 \uc778\uc0c1\uc758 \uba74\uc811\ub8e9 \ube44\uad50\ud558\uae30"',
            '\u2464 CTA: "\uc774 \ucf54\ub514\ub85c \uacb0\uc815" \u2192 \uc124\ubb38 \u2192 \uc644\ub8cc',
        ],
        "\uba74\uc811 \uc0c1\ud669\uc73c\ub85c \ub370\ubaa8\ub97c \ubcf4\uc5ec\ub4dc\ub9ac\uaca0\uc2b5\ub2c8\ub2e4.",
    ),
    (
        "MVP\ub294 \uac00\ubccd\uac8c, \uad6c\uc870\ub294 \ud655\uc7a5 \uac00\ub2a5\ud558\uac8c",
        '"\ube60\ub978 \uac80\uc99d\uacfc \ud655\uc7a5 \uac00\ub2a5\uc131\uc744 \ub3d9\uc2dc\uc5d0 \ud655\ubcf4\ud588\ub2e4"',
        [
            "\u2022 \ud604\uc7ac: JSON \ud30c\uc77c \uae30\ubc18 (outfits_scored.json, 1,645\uac1c \ucf54\ub514)",
            "\u2022 \uc774\uc720: MVP \ub2e8\uacc4\uc5d0\uc11c DB \uc138\ud305\ubcf4\ub2e4 \uc2e4\ud5d8 \uc18d\ub3c4 \uc6b0\uc120",
            "\u2022 \uc124\uacc4\ub41c \ubaa8\ub378: User / Outfit / Score / Reaction / StyleSeed",
            "\u2022 \ud655\uc7a5 \uacbd\ub85c: PostgreSQL(Supabase) \uc804\ud658 \uc900\ube44 \uc644\ub8cc",
            "\u2022 \ud589\ub3d9 \ub370\uc774\ud130(TTD, \uc2e0\ub8b0\ub3c4) \u2192 \uac1c\uc778\ud654 \ud559\uc2b5 \uae30\ubc18",
        ],
        "\ud604\uc7ac\ub294 \ube60\ub978 \uac80\uc99d\uc744 \uc704\ud574 JSON\uc744 \uc0ac\uc6a9\ud558\uc9c0\ub9cc, \uad6c\uc870\ub294 DB \ud655\uc7a5\uc744 \uace0\ub824\ud574 \uc124\uacc4\ud588\uc2b5\ub2c8\ub2e4.",
    ),
    (
        "\ud604\uc7ac \uace0\ubbfc: \uc804\ubb38\uac00 \uc218\uc900\uc758 \uc124\ub4dd\ub825",
        '"\ucd94\ucc9c\uc740 \ub418\uc9c0\ub9cc \uac10\ud0c4\uae4c\uc9c0\ub294 \uc544\uc9c1 \ubd80\uc871\ud558\ub2e4"',
        [
            '\u2022 \uc870\uac74 \uae30\ubc18 \ucd94\ucc9c\uc740 \ub3d9\uc791\ud558\uc9c0\ub9cc "\uc804\ubb38\uac00\uac00 \uace8\ub77c\uc900 \ub290\ub08c" \ubd80\uc871',
            '\u2022 \ucee8\ub7ec \uc870\ud569 \uc124\uba85: "\uc798 \uc5b4\uc6b8\ub824\uc694" \u2192 \uad6c\uccb4\uc801 \uadfc\uac70 \ud544\uc694',
            "\u2022 TPO\ubcc4 \uc2a4\ud0c0\uc77c \ud310\ub2e8 \uae30\uc900\uc774 \ub354 \uc815\uad50\ud574\uc838\uc57c \ud568",
            "\u2022 \ud604\uc7ac: Stylist Criteria Layer (\ud504\ub860\ud2b8 \ubb38\uad6c \ub808\ubca8)",
            "\u2022 \ub2e4\uc74c: \ucf54\ub514 \uc870\ud569 \uaddc\uce59 \uc790\uccb4\uc758 \uc804\ubb38\uc131 \uac15\ud654",
        ],
        "\uc194\uc9c1\ud558\uac8c \ub9d0\uc500\ub4dc\ub9ac\uba74, \uc0ac\uc6a9\uc790\uac00 \uc804\ubb38\uac00\uac00 \uace8\ub77c\uc92c\ub2e4\uace0 \uac10\ud0c4\ud558\ub294 \uc218\uc900\uc740 \uc544\uc9c1 \uc544\ub2d9\ub2c8\ub2e4.",
    ),
    (
        "\uba58\ud1a0 \ud53c\ub4dc\ubc31 \uc694\uccad",
        '"\ubc29\ud5a5\uc744 \uc7a1\uc558\uace0, \ub2e4\uc74c \ub2e8\uacc4\uc758 \uae4a\uc774\ub97c \uacb0\uc815\ud574\uc57c \ud569\ub2c8\ub2e4"',
        [
            "\u2022 \ud604\uc7ac \ubc29\ud5a5: Decision-First + Guided Decision UX",
            "\u2022 \uc9c8\ubb38 1: \uc804\ubb38\uac00 \ud310\ub2e8 \uae30\uc900\uc744 \uc2dc\uc2a4\ud15c\ud654\ud558\ub294 \ucd5c\uc801 \ubc29\ubc95\uc740?",
            "\u2022 \uc9c8\ubb38 2: MVP\uc5d0\uc11c \uac1c\uc778\ud654 \uc218\uc900\uc744 \uc5b4\ub514\uae4c\uc9c0 \uac00\uc838\uac00\uc57c \ud558\ub294\uac00?",
            "",
            "\uac10\uc0ac\ud569\ub2c8\ub2e4.",
        ],
        "ColorFit\uc758 \ubc29\ud5a5\uc740 Decision \uc11c\ube44\uc2a4\uc785\ub2c8\ub2e4. \uba58\ud1a0\ub2d8\uaed8 \uc5ec\ucad1\uace0 \uc2f6\uc740 \uac74 \uc804\ubb38\uac00 \ud310\ub2e8\uc744 \uad6c\uc870\ud654\ud558\ub294 \ubc29\ubc95\uacfc MVP\uc5d0\uc11c \uac1c\uc778\ud654 \ubc94\uc704\uc785\ub2c8\ub2e4.",
    ),
]

for i, (t, sub, bs, note) in enumerate(SLIDES):
    make_slide(t, sub, bs, note, str(i + 1))

make_qa()

prs.save("ColorFit_멘토발표.pptx")
print("PPT 생성 완료: ColorFit_멘토발표.pptx")
